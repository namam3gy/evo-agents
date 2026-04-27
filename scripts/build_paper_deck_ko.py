"""Build PPT 2 — paper-style deck (Korean, ~25 slides, lean), Forest & Moss palette.

Outputs to docs/presentations/2026-04-26-evo-agents/2-paper-style-ko.pptx.
Run via:
    uv run python scripts/build_paper_deck_ko.py
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

REPO_ROOT = Path(__file__).resolve().parents[1]
DECK_DIR = REPO_ROOT / "docs" / "presentations" / "2026-04-26-evo-agents"
ASSETS = DECK_DIR / "assets"
OUT_PATH = DECK_DIR / "2-paper-style-ko.pptx"

C_FOREST = RGBColor(0x2C, 0x5F, 0x2D)
C_MOSS = RGBColor(0x97, 0xBC, 0x62)
C_CREAM = RGBColor(0xF5, 0xF5, 0xF5)
C_CHARCOAL = RGBColor(0x36, 0x45, 0x4F)
C_ACCENT = RGBColor(0xB8, 0x50, 0x42)
C_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
C_LIGHT_GREY = RGBColor(0xCF, 0xD3, 0xD7)

# Korean-friendly fonts. NanumGothic + Noto Sans CJK are widely available on Linux/macOS;
# Windows PowerPoint will fall back to Malgun Gothic via system substitution.
FONT_HEADER = "NanumGothic"
FONT_BODY = "NanumGothic"

SLIDE_W_IN = 13.333
SLIDE_H_IN = 7.5

# ---------------------------------------------------------------------------
# helpers (mirrored from build_advisor_deck with Korean fonts)
# ---------------------------------------------------------------------------


def _set_slide_background(slide, color: RGBColor) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_textbox(slide, x_in, y_in, w_in, h_in, text,
                 *, size=14, color=C_CHARCOAL, bold=False, italic=False,
                 font=FONT_BODY, align=PP_ALIGN.LEFT,
                 anchor=MSO_ANCHOR.TOP, line_spacing=1.2):
    tx = slide.shapes.add_textbox(Inches(x_in), Inches(y_in),
                                  Inches(w_in), Inches(h_in))
    tf = tx.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Inches(0.06)
    tf.margin_right = Inches(0.06)
    tf.margin_top = Inches(0.04)
    tf.margin_bottom = Inches(0.04)
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = line_spacing
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    r.font.name = font
    return tx


def _add_bullets(slide, x_in, y_in, w_in, h_in, items,
                 *, size=15, color=C_CHARCOAL, bullet="•",
                 font=FONT_BODY, line_spacing=1.3, sub_size=13):
    tx = slide.shapes.add_textbox(Inches(x_in), Inches(y_in),
                                  Inches(w_in), Inches(h_in))
    tf = tx.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.06)
    tf.margin_right = Inches(0.06)
    tf.margin_top = Inches(0.04)
    tf.margin_bottom = Inches(0.04)
    for i, raw in enumerate(items):
        if isinstance(raw, str):
            text, lvl, run_color = raw, 0, color
        else:
            text = raw[0]
            lvl = raw[1] if len(raw) > 1 else 0
            run_color = raw[2] if len(raw) > 2 else color
        prefix = "  " * lvl + (bullet if lvl == 0 else "–") + " "
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = line_spacing
        r = p.add_run()
        r.text = prefix + text
        r.font.size = Pt(size if lvl == 0 else sub_size)
        r.font.color.rgb = run_color
        r.font.name = font
    return tx


def _add_title(slide, title_text, *, color=C_CHARCOAL, size=30, x_in=0.6, y_in=0.55,
               w_in=12.1, h_in=0.9):
    return _add_textbox(slide, x_in, y_in, w_in, h_in, title_text,
                        size=size, color=color, bold=True, font=FONT_HEADER,
                        align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)


def _add_accent_strip(slide, color=C_FOREST, height_in=0.18):
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0,
                                  Inches(SLIDE_W_IN), Inches(height_in))
    rect.fill.solid()
    rect.fill.fore_color.rgb = color
    rect.line.fill.background()


def _add_image(slide, path: Path, x_in, y_in, w_in, h_in=None):
    if h_in is None:
        return slide.shapes.add_picture(str(path), Inches(x_in), Inches(y_in),
                                        width=Inches(w_in))
    return slide.shapes.add_picture(str(path), Inches(x_in), Inches(y_in),
                                    width=Inches(w_in), height=Inches(h_in))


def _add_pill(slide, x_in, y_in, w_in, h_in, text,
              *, fill=C_MOSS, font_color=C_WHITE, size=12, bold=True):
    rect = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  Inches(x_in), Inches(y_in),
                                  Inches(w_in), Inches(h_in))
    rect.fill.solid()
    rect.fill.fore_color.rgb = fill
    rect.line.fill.background()
    tf = rect.text_frame
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.03)
    tf.margin_bottom = Inches(0.03)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = font_color
    r.font.name = FONT_HEADER
    return rect


def _add_table(slide, x_in, y_in, w_in, h_in, data, *,
               header_fill=C_FOREST, header_font=C_WHITE,
               body_font=C_CHARCOAL, body_size=12,
               header_size=13, alt_fill=C_CREAM):
    rows = len(data)
    cols = len(data[0])
    table = slide.shapes.add_table(rows, cols,
                                   Inches(x_in), Inches(y_in),
                                   Inches(w_in), Inches(h_in)).table
    for j, cell_text in enumerate(data[0]):
        cell = table.cell(0, j)
        cell.fill.solid()
        cell.fill.fore_color.rgb = header_fill
        tf = cell.text_frame
        tf.margin_left = Inches(0.05)
        tf.margin_right = Inches(0.05)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = str(cell_text)
        r.font.size = Pt(header_size)
        r.font.bold = True
        r.font.color.rgb = header_font
        r.font.name = FONT_HEADER
    for i in range(1, rows):
        row_alt = alt_fill if i % 2 else C_WHITE
        for j, cell_text in enumerate(data[i]):
            cell = table.cell(i, j)
            cell.fill.solid()
            cell.fill.fore_color.rgb = row_alt
            tf = cell.text_frame
            tf.margin_left = Inches(0.06)
            tf.margin_right = Inches(0.06)
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER if j > 0 else PP_ALIGN.LEFT
            r = p.add_run()
            r.text = str(cell_text)
            r.font.size = Pt(body_size)
            r.font.color.rgb = body_font
            r.font.name = FONT_BODY
    return table


def make_title_slide(prs, title, subtitle, footer):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_background(slide, C_FOREST)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                 Inches(0), Inches(3.0),
                                 Inches(0.3), Inches(1.5))
    bar.fill.solid()
    bar.fill.fore_color.rgb = C_MOSS
    bar.line.fill.background()
    _add_textbox(slide, 0.9, 2.4, 12.0, 1.6, title,
                 size=38, color=C_WHITE, bold=True, font=FONT_HEADER,
                 anchor=MSO_ANCHOR.MIDDLE)
    _add_textbox(slide, 0.9, 4.1, 12.0, 1.0, subtitle,
                 size=18, color=C_LIGHT_GREY, font=FONT_BODY, italic=True,
                 anchor=MSO_ANCHOR.TOP)
    _add_textbox(slide, 0.9, 6.7, 12.0, 0.5, footer,
                 size=12, color=C_LIGHT_GREY, font=FONT_BODY,
                 anchor=MSO_ANCHOR.TOP)
    return slide


def make_section_slide(prs, label, body_text=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_background(slide, C_FOREST)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                 Inches(0.6), Inches(3.4),
                                 Inches(0.25), Inches(0.9))
    bar.fill.solid()
    bar.fill.fore_color.rgb = C_MOSS
    bar.line.fill.background()
    _add_textbox(slide, 1.1, 3.3, 12.0, 1.0, label,
                 size=38, color=C_WHITE, bold=True, font=FONT_HEADER,
                 anchor=MSO_ANCHOR.MIDDLE)
    if body_text:
        _add_textbox(slide, 1.1, 4.4, 12.0, 1.4, body_text,
                     size=16, color=C_LIGHT_GREY, font=FONT_BODY, italic=True,
                     anchor=MSO_ANCHOR.TOP)
    return slide


def make_content_slide(prs, title, *, kicker=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_background(slide, C_CREAM)
    _add_accent_strip(slide, color=C_FOREST)
    if kicker:
        _add_textbox(slide, 0.6, 0.27, 12.0, 0.32, kicker,
                     size=11, color=C_MOSS, bold=True, font=FONT_HEADER,
                     anchor=MSO_ANCHOR.TOP)
    _add_title(slide, title, y_in=0.55)
    return slide


def make_conclusion_slide(prs, big, sub):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_background(slide, C_FOREST)
    _add_textbox(slide, 0.6, 2.6, 12.1, 1.3, big,
                 size=46, color=C_WHITE, bold=True, font=FONT_HEADER,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    _add_textbox(slide, 0.6, 4.1, 12.1, 1.6, sub,
                 size=18, color=C_LIGHT_GREY, font=FONT_BODY,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP)
    return slide


# ---------------------------------------------------------------------------
# slides
# ---------------------------------------------------------------------------


def build():
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W_IN)
    prs.slide_height = Inches(SLIDE_H_IN)

    # ---- Title (1) -----------------------------------------------------
    make_title_slide(
        prs,
        "성찰만으로 진화하는 멀티-에이전트 DAG",
        "검색·강화학습 없이 in-context reflection 만으로 그래프를 점진적으로 편집하는 파일럿",
        "thyun.park · 논문화 형식 발표 자료 · 2026-04-27 (v3 반영본)",
    )

    # ---- Abstract (2) --------------------------------------------------
    s = make_content_slide(prs, "초록 (Abstract)")
    _add_textbox(s, 0.6, 1.4, 12.1, 5.7,
                 "본 연구는 멀티-에이전트 DAG (위상 + 페르소나 + 엣지)를 LLM 컨트롤러가 "
                 "오직 in-context 성찰(reflection)만으로 점진적 편집할 수 있는지 묻는 파일럿 "
                 "연구이다. 컨트롤러는 검색 알고리즘이나 RL 학습 없이, 워커 에이전트들의 "
                 "trajectory tape을 읽어 다음 라운드에 적용할 EditBatch만 출력한다.\n\n"
                 "Qwen2.5-32B 백본 위에서 FinanceBench / MEDIQ / AgentClinic 세 도메인을 "
                 "대상으로 CoT, Planner-Executor 두 베이스라인과 진화 그래프를 비교한다.\n\n"
                 "v1은 'verifier 추가' 반사만 보였고, v2는 'organization-designer' 프레이밍 + "
                 "도메인 brief로 전문가 페르소나(예: differential_diagnostician, gaap_analyst)를 "
                 "출력. 측정 노이즈 ±13pp 문제를 해결하기 위한 streaming paired-batch 모드는 "
                 "두 시드에서 Δ vs P-E = +4pp로 재현. 그러나 streaming의 cross-batch absolute "
                 "acc, max_agents binding, orphan-edit INVALID 라는 구조적 한계를 보고 "
                 "Controller v3 (full-pass + 샘플 단위 reflection 30→3→1 hierarchical "
                 "aggregation + 워커 DAG의 [SUMMARY]/[QUERY] side-channel Q&A + orphan "
                 "auto-drop)를 설계·구현. v3 첫 run (MEDIQ s=0)에서 Evolved 52% > CoT 46% "
                 "(+6pp) 로 파일럿 최초 CoT-beat 결과 (단일 시드). multi-seed sweep이 "
                 "venue submission 직전 잔여 작업.",
                 size=13, color=C_CHARCOAL, line_spacing=1.35)

    # ---- Section: Introduction ----------------------------------------
    make_section_slide(prs, "1. 서론",
                       "동기 · 문제 · 기여")

    # ---- Intro (3, 4) --------------------------------------------------
    s = make_content_slide(prs, "1.1 연구 동기", kicker="서론")
    _add_bullets(s, 0.6, 1.4, 12.1, 5.6, [
        "단일 LLM은 강하지만, 도메인-특화 작업에서 'planner / specialist / verifier' 식 "
        "역할 분리가 정확도를 끌어올린다는 보고가 다수.",
        "그러나 멀티-에이전트 그래프(누가 누구에게 input을 주는지)를 어떻게 자동으로 "
        "디자인할지는 여전히 열린 문제.",
        "기존 접근은 (a) 검색 — ADAS / AFlow / GPTSwarm / MaAS / Puppeteer (NeurIPS 2025), "
        "(b) RL 미세조정 두 갈래.",
        ("질문: 학습이나 검색 없이, 단지 trajectory tape을 LLM에게 보여주고 그래프를 "
         "직접 편집하라고 하면 (= 성찰), 그것만으로 진화가 일어날 수 있는가?", 0, C_FOREST),
    ], size=15)

    s = make_content_slide(prs, "1.2 기여 (Contributions)", kicker="서론")
    _add_pill(s, 0.6, 1.3, 12.1, 0.4,
              "(1) Reflection-only 멀티-에이전트 진화 — 검색·RL 없이 그래프를 점진적 편집",
              fill=C_FOREST, size=12)
    _add_textbox(s, 0.6, 1.72, 12.1, 0.6,
                 "Trajectory tape → 컨트롤러 LLM → EditBatch → apply_edits() 의 단순 루프.",
                 size=11, color=C_CHARCOAL)
    _add_pill(s, 0.6, 2.42, 12.1, 0.4,
              "(2) v2 컨트롤러: organization-designer 프레이밍 + 도메인 brief",
              fill=C_FOREST, size=12)
    _add_textbox(s, 0.6, 2.84, 12.1, 0.6,
                 "도메인 어휘 강제 + 전문가 persona 규칙으로 generic 'verifier' 반사 차단; "
                 "같은 backbone이 도메인별로 다른 전문가처럼 행동.",
                 size=11, color=C_CHARCOAL)
    _add_pill(s, 0.6, 3.54, 12.1, 0.4,
              "(3) Streaming paired-batch — paired ACCEPT 1/9 → 4/10, Δ vs P-E = +4pp 두 시드 재현",
              fill=C_FOREST, size=12)
    _add_textbox(s, 0.6, 3.96, 12.1, 0.6,
                 "best vs candidate를 SAME batch로 평가해 cross-batch 노이즈 cancel.",
                 size=11, color=C_CHARCOAL)
    _add_pill(s, 0.6, 4.66, 12.1, 0.4,
              "(4) Controller v3 — full-pass + 샘플 단위 reflection + side-channel Q&A",
              fill=C_ACCENT, size=12, font_color=C_WHITE)
    _add_textbox(s, 0.6, 5.08, 12.1, 1.95,
                 "(a) train n=30 sample-level evals → 30→3→1 hierarchical aggregation; "
                 "(b) 워커 DAG에 [SUMMARY] transcript + [QUERY] one-shot Q&A; "
                 "(c) orphan auto-drop으로 INVALID 0; (d) MEDIQ s=0에서 Evolved 52% > CoT 46% "
                 "(+6pp) — 파일럿 최초 CoT-beat 결과 (단일 시드).",
                 size=11, color=C_CHARCOAL)

    # ---- Section: Related Work -----------------------------------------
    make_section_slide(prs, "2. 관련 연구",
                       "검색-기반 · RL-기반 · 본 연구의 자리")

    # ---- Related work (5, 6) -------------------------------------------
    s = make_content_slide(prs, "2.1 검색·RL 기반 멀티-에이전트 자동 설계", kicker="관련 연구")
    _add_table(s, 0.6, 1.5, 12.1, 4.4, [
        ["방법", "메커니즘", "특징", "본 연구 대비"],
        ["ADAS (2024)", "Archive 검색", "메타-에이전트가 candidate archive에서 선택", "검색 archive 사용 — 본 연구는 미사용"],
        ["AFlow", "MCTS", "워크플로우 트리 탐색", "MCTS 사용 — 본 연구는 미사용"],
        ["GPTSwarm", "Graph search", "Edge importance 학습", "학습 신호 사용 — 본 연구는 미사용"],
        ["MaAS", "Supernet", "Multi-agent supernet 가지치기", "supernet 학습 사용 — 본 연구는 미사용"],
        ["Puppeteer (NeurIPS '25)", "RL", "마스터 에이전트가 액션-같은 협업 결정", "RL 사용 — 본 연구는 미사용"],
    ])
    _add_textbox(s, 0.6, 6.2, 12.1, 0.7,
                 "전부 검색 절차나 학습 신호를 가정. 본 연구는 frozen backbone + in-context reflection만 사용.",
                 size=13, italic=True, color=C_FOREST)

    s = make_content_slide(prs, "2.2 본 연구의 차별점", kicker="관련 연구")
    _add_bullets(s, 0.6, 1.4, 12.1, 5.6, [
        ("본 연구는 frozen LLM + 자연어 trajectory tape + EditBatch 출력만으로 그래프를 "
         "co-evolve (위상 + 페르소나 + 엣지) 한다. 학습도, 검색도 없다.", 0, C_FOREST),
        "그래서 'reflection-only' 라는 좁은 슬롯에 위치한다 — 좁다는 것은 곧 reviewer가 "
        "쉽게 공격할 수 있는 부분이기도 하다.",
        "Caveat (project.md 기준): 이 차별점만으로는 NeurIPS 2026 main 합격 기준을 자체 "
        "평가에서 충족하지 않는다 (3-6% acceptance 추정). EMNLP 2026 ARR (~5-10%)이 현실적 1차 목표.",
        "약점에 대한 정직한 정리는 §10 (한계)에서 다룬다.",
    ], size=14)

    # ---- Section: Problem & Definitions -------------------------------
    make_section_slide(prs, "3. 문제 정의",
                       "DAG / Agent / Edit · Reflection-only 가정")

    # ---- Problem (7, 8) -----------------------------------------------
    s = make_content_slide(prs, "3.1 그래프와 편집 연산", kicker="문제 정의")
    _add_image(s, ASSETS / "dag_baselines.png", 0.4, 1.4, w_in=12.5)
    _add_bullets(s, 0.6, 5.2, 12.1, 1.9, [
        "그래프 G = (Agents A, Edges E ⊆ A × A), START / END 예약 노드 포함; orchestrator가 "
        "topological order로 실행.",
        "Agent = (name, persona text). Persona는 해당 노드에서 LLM이 조건화되는 텍스트.",
        ("편집 연산: add_agent / remove_agent / rewrite_persona / add_edge / remove_edge. "
         "한 라운드에 1-3 연산을 묶어 EditBatch로 출력.", 0, C_CHARCOAL),
    ], size=13)

    s = make_content_slide(prs, "3.2 Reflection-only 가정", kicker="문제 정의")
    _add_bullets(s, 0.6, 1.4, 12.1, 5.6, [
        "백본 weights는 동결 (no fine-tune, no RL).",
        "컨트롤러 LLM은 검색 archive를 가지지 않는다 — 받는 입력은 trajectory tape + (필요시) "
        "domain brief뿐.",
        "출력 = 다음 그래프 상태에 적용할 EditBatch (rationale + edits 리스트).",
        "Accept rule = (legacy) val_acc 비교 / (streaming) paired b_acc vs c_acc on the same batch.",
        ("이 가정은 의도적으로 좁다. 'in-context reflection 만으로 어디까지 가능한가?'를 "
         "측정하기 위해.", 0, C_FOREST),
    ], size=14)

    # ---- Section: Method ----------------------------------------------
    make_section_slide(prs, "4. 방법 (Method)",
                       "Controller v2 · Domain brief · Streaming paired-batch")

    # ---- Method (9, 10, 11, 12) ---------------------------------------
    s = make_content_slide(prs, "4.1 Controller v2 — organization designer 프레이밍",
                           kicker="방법")
    _add_image(s, ASSETS / "controller_v1_v2_loop.png", 0.4, 1.4, w_in=12.5)
    _add_textbox(s, 0.6, 6.4, 12.1, 0.6,
                 "v1과 v2는 같은 reflection 루프를 공유한다. 차이는 SYSTEM prompt와 "
                 "side-channel domain brief뿐.",
                 size=12, italic=True, color=C_MOSS)

    s = make_content_slide(prs, "4.2 Persona 작성 규칙 + 도메인 brief", kicker="방법")
    _add_bullets(s, 0.6, 1.4, 12.1, 5.6, [
        "각 신규 agent는 (a) 인용된 specialty + (b) 구체적 procedure를 페르소나 텍스트에 포함해야 함.",
        "Generic 'verifier / summarizer / critic'는 specialty와 짝짓지 않으면 금지.",
        "도메인 brief (data/briefs/{financebench,mediq,agentclinic}.md, ~80-110 줄):",
        ("task style, failure modes, useful expertise, useful patterns, anti-patterns 등을 "
         "포함한 사전 작성된 자연어 문서.", 1),
        ("브리프는 user prompt에 prefix로 포함되어 컨트롤러가 같은 backbone임에도 "
         "도메인별로 다른 전문가처럼 행동하게 만든다.", 0, C_FOREST),
        "Anti-repeat rule: '거부된 edit을 반복하지 말 것' + 'rename ≠ different idea' (concept-level).",
    ], size=14)

    s = make_content_slide(prs, "4.3 Streaming paired-batch evolve mode", kicker="방법")
    _add_image(s, ASSETS / "paired_batch_flow.png", 0.4, 1.4, w_in=12.5)
    _add_bullets(s, 0.6, 5.7, 12.1, 1.5, [
        "라운드별 부트스트랩 미니배치 (B 태스크, train+val 풀에서 with-replacement) 추출.",
        "best_graph와 candidate (apply_edits 결과) 둘 다 같은 batch로 평가 → b_acc, c_acc.",
        "ACCEPT iff c_acc > b_acc + ε. Cross-batch 난이도 drift는 paired 비교에서 cancel.",
    ], size=13)

    # ---- §4.4 Controller v3 motivation + mode (NEW) -----------------
    s = make_content_slide(prs, "4.4 Controller v3 — full-pass + strict accept",
                           kicker="방법")
    _add_pill(s, 0.6, 1.3, 12.1, 0.45,
              "Streaming run #1·#2가 노출한 3대 문제",
              fill=C_ACCENT, size=12, font_color=C_WHITE)
    _add_bullets(s, 0.6, 1.78, 12.1, 1.85, [
        ("(P1) cross-batch best_val_acc 비신뢰 — max over 독립 부트스트랩 batch.", 0, C_ACCENT),
        ("(P2) max_agents=6 binding — run #1·#2 6 round을 INVALID로 잃음.", 0, C_ACCENT),
        ("(P3) orphan-edit INVALID — partial-wire add_agent / orphan prune.", 0, C_ACCENT),
    ], size=12)
    _add_pill(s, 0.6, 3.7, 12.1, 0.45,
              "v3의 응답",
              fill=C_FOREST, size=12)
    _add_bullets(s, 0.6, 4.18, 12.1, 2.9, [
        ("Mode: streaming → legacy-style full-pass per iter (n_train=30 동일 task 재사용).", 0, C_CHARCOAL),
        ("Accept rule: train_acc(candidate) > train_acc(best) — strict.", 0, C_CHARCOAL),
        ("max_agents: 10 HARD (apply_edits가 INVALID), max_edges: 50 soft (prompt only).", 0, C_CHARCOAL),
        ("apply_edits가 무입력 / 무출력 agent를 silently drop → partial-wire 편집이 INVALID 대신 no-op.", 0, C_FOREST),
        ("모드 자체는 wall ½ (10h → 5h on MEDIQ s=0); reflection density는 §4.5/§4.6에서 강화.", 0, C_FOREST),
    ], size=12)

    # ---- §4.5 Sample-level + hierarchical aggregation (NEW) ----------
    s = make_content_slide(prs, "4.5 샘플 단위 reflection + 30 → 3 → 1 계층적 집계",
                           kicker="방법")
    _add_textbox(s, 0.6, 1.3, 12.1, 0.5,
                 "각 train task 마다 컨트롤러를 1회 호출 → 그 한 샘플의 tape이 그래프에 대해 무엇을 시사하는지 평가.",
                 size=12, italic=True, color=C_CHARCOAL)
    _add_pill(s, 0.6, 1.85, 5.8, 0.45, "eval_sample 출력 스키마", fill=C_FOREST, size=12)
    _add_bullets(s, 0.6, 2.32, 5.8, 4.7, [
        "rationale (2-4 문장, tape 인용)",
        "suggested_edits (0-3개 candidate)",
        "priority: 0-100 (자기 의견 강도)",
        "target_aspect ∈ {structure, role, length, expertise}",
        ("→ 컨트롤러가 자가-보정을 통해 노이즈한 샘플에 가중치를 줄임.", 0, C_MOSS),
    ], size=12)
    _add_pill(s, 6.7, 1.85, 6.0, 0.45, "Hierarchical aggregation",
              fill=C_FOREST, size=12)
    _add_bullets(s, 6.7, 2.32, 6.0, 4.7, [
        "30 sample-evals (priority 내림차순)",
        "  → groups of 10 → 3 mid_decisions",
        "  → 1 final EditBatch → apply_edits",
        ("aggregate_mid: 패턴 공유 샘플에 가중, priority=100 outlier에 과반응 금지.", 1),
        ("이 구조가 한 샘플의 노이즈가 그래프 편집을 좌우하지 못하게 막음.", 0, C_FOREST),
    ], size=12)

    # ---- §4.6 Side-channel Q&A on worker DAG (NEW) -------------------
    s = make_content_slide(prs, "4.6 [SUMMARY] / [QUERY] — 워커 DAG side-channel",
                           kicker="방법")
    _add_pill(s, 0.6, 1.3, 5.8, 0.45,
              "[SUMMARY] (S-2) — agent 출력 강제 구조",
              fill=C_FOREST, size=11)
    _add_bullets(s, 0.6, 1.78, 5.8, 5.3, [
        "각 워커는 응답 끝에 다음 블록을 emit:",
        ("  [SUMMARY] claim / evidence / confidence [/SUMMARY]", 1),
        ("orchestrator가 정규식으로 추출 → 다음 워커의 prompt 앞에 [Conversation so far] "
         "블록으로 prepend.", 0, C_CHARCOAL),
        ("→ agent.inputs 와는 *독립된* 추가 정보 채널 (W-2: 같은 task에 대해 모든 워커가 "
         "claim/evidence/confidence를 공유).", 0, C_FOREST),
        ("Fallback: 누락 시 마지막 문장으로 claim 채움 (no INVALID).", 0, C_CHARCOAL),
    ], size=11)
    _add_pill(s, 6.7, 1.3, 6.0, 0.45,
              "[QUERY] (Q-3) — one-shot peer Q&A",
              fill=C_FOREST, size=11)
    _add_bullets(s, 6.7, 1.78, 6.0, 5.3, [
        ("워커 prompt에 \"원하면 다른 에이전트 1명에게 [QUERY <name>] <질문> 한 줄 출력 후 "
         "STOP\" instruction 추가.", 0, C_CHARCOAL),
        ("Call A (메인) → [QUERY] 검출 시 Call B (lite-mode 답변, 256 tok, 재귀 금지) → "
         "Call C (메인 재개, 답변 부착)", 0, C_CHARCOAL),
        ("Hard rules: 답변자는 query 금지 / 1 query/노드/태스크 / 노드당 최악 3 LLM 호출.", 0, C_CHARCOAL),
        ("→ 워커 단위 reflection density 증가, 단순 cascade가 아닌 협업 가능.", 0, C_FOREST),
    ], size=11)

    s = make_content_slide(prs, "4.7 Accept rule + max_agents 제약 (3 모드 비교)", kicker="방법")
    _add_table(s, 0.6, 1.4, 12.1, 2.4, [
        ["모드", "Accept rule", "max_agents", "기타 제약"],
        ["Legacy (v1, v2 n=30)", "val_acc > best (strict)", "8 (soft)", "tie REJECT"],
        ["Streaming (v2 §5.1)", "paired Δ > ε on same batch", "8 (soft)", "concept-level repeat (soft)"],
        ["v3 full-pass (§5.2)", "train_acc(cand) > train_acc(best)", "10 HARD", "orphan auto-drop, max_edges=50 soft"],
    ])
    _add_bullets(s, 0.6, 4.1, 12.1, 2.9, [
        ("v3에서 max_agents가 HARD인 이유: prompt-only soft가 컨트롤러에 무시되어 streaming "
         "30~40% INVALID를 만든 사례 (run #1·#2). HARD로 옮기되 cap을 8 → 10으로 완화.", 0, C_FOREST),
        ("orphan auto-drop: apply_edits가 incoming/outgoing 0인 노드를 조용히 제거 → "
         "v3 첫 run에서 INVALID 0/10 달성 (streaming 4/10·2/10 대비).", 0, C_FOREST),
        ("3 모드 모두 reflection-only (학습 없음, 검색 archive 없음).", 0, C_CHARCOAL),
    ], size=12)

    # ---- Implementation (13) ------------------------------------------
    s = make_content_slide(prs, "5. 구현 (Implementation)", kicker="구현")
    _add_bullets(s, 0.6, 1.4, 6.0, 5.6, [
        "src/llm.py — OpenAI-compatible chat client (vLLM @ :8000)",
        "src/graph.py — Graph + edit ops + DAG validity 검사",
        "src/orchestrator.py — Task → Tape 실행",
        "src/controller.py — Reflection LLM, propose_edits()",
        "src/evolve.py — legacy + streaming evolve loops",
        "src/score.py — MCQ exact-match (MEDIQ) + LLM-judge",
        "scripts/run_pilot.py — 드라이버",
        "scripts/serve_vllm.sh — vLLM 0.19.1 + gcc 자동 설치",
    ], size=13)
    _add_pill(s, 6.9, 1.5, 5.8, 0.5, "백본: Qwen2.5-32B-Instruct (bf16, single H200)",
              fill=C_FOREST, size=12)
    _add_bullets(s, 6.9, 2.1, 5.8, 4.9, [
        "vLLM + prefix caching",
        "EVO_GPU_UTIL=0.55 (공유 GPU)",
        "max_model_len = 16384",
        "uv-managed env, .python-version pin",
        "각 run: results/<run_id>/{evolve_log.json, results.json, plots/}",
        ("코드 + 결과: github.com/namam3gy (project repo)", 0, C_MOSS),
    ], size=12)

    # ---- Experimental Setup (14, 15) ----------------------------------
    s = make_content_slide(prs, "6.1 실험 설정 — 도메인", kicker="실험 설정")
    _add_pill(s, 0.6, 1.4, 4.0, 0.5, "FinanceBench (LLM-judge)", fill=C_FOREST, size=13)
    _add_bullets(s, 0.6, 1.95, 4.0, 5.0, [
        "10-K 추출 QA",
        "GAAP / TTM / 회계연도",
        "긴 컨텍스트 retrieval",
        "실패: hedging, 기간 혼동",
    ], size=12)
    _add_pill(s, 4.7, 1.4, 4.0, 0.5, "MEDIQ (MCQ exact-match)", fill=C_FOREST, size=13)
    _add_bullets(s, 4.7, 1.95, 4.0, 5.0, [
        "임상 vignette → MCQ",
        "감별 진단 (differential dx)",
        "Base-rate reasoning",
        "실패: 인구학 무시",
    ], size=12)
    _add_pill(s, 8.8, 1.4, 4.0, 0.5, "AgentClinic (LLM-judge)", fill=C_FOREST, size=13)
    _add_bullets(s, 8.8, 1.95, 4.0, 5.0, [
        "single-pass 임상 케이스",
        "Specialty routing",
        "Decisive answer 요구",
        "실패: triage 오류",
    ], size=12)
    _add_textbox(s, 0.6, 6.5, 12.1, 0.6,
                 "공통: Qwen2.5-32B-Instruct, val/test sample 분할 별도, seed 명시.",
                 size=12, italic=True, color=C_CHARCOAL)

    s = make_content_slide(prs, "6.2 측정 지표 + 노이즈 floor", kicker="실험 설정")
    _add_bullets(s, 0.6, 1.4, 12.1, 5.6, [
        "Test acc — 보유한 test split (n=30 또는 n=50)에서의 정확도.",
        "Paired-accept rate — streaming mode에서 ACCEPT된 라운드 수 / 전체 유효 라운드.",
        "Token cost — controller_tokens + worker_tokens 합 (run 단위).",
        ("Same-graph 노이즈 floor ≈ ±13pp at n=30 (FinanceBench v2 retry 사례; vLLM batch-ordering, "
         "KV-cache state 비결정성, sample 분할 조합).", 0, C_ACCENT),
        ("→ n=30 단일 시드에서 도출한 architectural Δ는 노이즈 안에 묻힘. "
         "다중 시드/배치 평균이 필수.", 0, C_FOREST),
    ], size=14)

    # ---- Section: Experiments -----------------------------------------
    make_section_slide(prs, "7. 실험 결과",
                       "calib_01 → v1 n=30 → v2 n=30 → streaming run #1 → patches sanity → §5.2 ongoing")

    # ---- Experiments (16-21, 6 slides) --------------------------------
    s = make_content_slide(prs, "7.1 calib_01 (GSM8K) — pilot의 출발점", kicker="실험 결과")
    _add_image(s, REPO_ROOT / "results" / "calib_01" / "plots" / "accuracy_vs_iter.png",
               0.6, 1.4, w_in=6.4)
    _add_bullets(s, 7.3, 1.5, 5.6, 5.6, [
        "Setup: n_val = n_test = 50, max_iters = 3, wall ~66 min.",
        "결과: Evolved 86 % / CoT 92 % / P-E 90 %.",
        ("Evolved가 두 baseline 모두에 패배 (test) + 2-3.7× 토큰 비용.", 0, C_ACCENT),
        ("Insight: GSM8K은 단일-모델 saturation + self-contained 텍스트 + 선형 산술 — "
         "정보 비대칭이 없어서 멀티-에이전트의 lever가 없는 도메인.", 0, C_FOREST),
        "→ 도메인 피봇 결정 (FinanceBench / MEDIQ / AgentClinic).",
    ], size=12)

    s = make_content_slide(prs, "7.2 v1 controller @ n=30 — 'verifier 반사' 노출", kicker="실험 결과")
    _add_table(s, 0.6, 1.5, 12.1, 2.4, [
        ["도메인", "CoT test", "P-E test", "Evolved test", "Δ vs best baseline"],
        ["FinanceBench", "70.0 %", "66.7 %", "70.0 %", "0 pp"],
        ["MEDIQ", "43.3 %", "43.3 %", "50.0 %", "+6.7 pp (±18 pp 노이즈)"],
        ["AgentClinic", "66.7 %", "70.0 %", "70.0 %", "0 pp"],
    ])
    _add_bullets(s, 0.6, 4.2, 12.1, 2.9, [
        ("FinanceBench: 같은 add_verifier edit을 3 라운드 연속 emit (도메인 어휘 0).", 0, C_ACCENT),
        ("MEDIQ: iter-2 ACCEPT 후 edit이 약간 다양화.", 0, C_CHARCOAL),
        ("AgentClinic: summarizer ↔ verifier 교대.", 0, C_CHARCOAL),
        ("Read: H1 약화 (3 도메인 모두 baseline 못 이김); v1은 도메인 헤드룸을 못 살림. "
         "v2 재설계로 이행.", 0, C_FOREST),
    ], size=13)

    s = make_content_slide(prs, "7.3 Controller v2 @ n=30 — 행동 변화 ↔ test 승리는 아님",
                           kicker="실험 결과")
    _add_table(s, 0.6, 1.5, 12.1, 2.4, [
        ["도메인", "CoT test", "P-E test", "Evolved test", "Δ", "best_graph"],
        ["FinanceBench (16k)", "73.3 %", "70.0 %", "83.3 %*", "+10 pp*", "seed (모두 REJECT)"],
        ["MEDIQ", "43.3 %", "46.7 %", "43.3 %", "−3.4 pp", "3-agent (iter-2 ACCEPT)"],
        ["AgentClinic", "60.0 %", "73.3 %", "66.7 %", "−6.6 pp", "seed (모두 REJECT)"],
    ])
    _add_bullets(s, 0.6, 4.2, 12.1, 2.9, [
        ("* FinanceBench의 +10 pp는 SAME-GRAPH 노이즈 — best_graph가 seed인데 test가 P-E 동시 "
         "런 (70%)과 13 pp 떨어짐. v2 승리로 인용 불가.", 0, C_ACCENT),
        ("Persona 이름이 도메인-특화로 변모: gaap_analyst, period_validator, "
         "differential_diagnostician, adolescent_specialist 등.", 0, C_FOREST),
        ("AgentClinic iter-3 (REJECTED): triage → {gastro, cardio} → END 형태의 정식 specialty "
         "department 구조 제안 — H2 BEHAVIOR 확인, TEST WIN 미달성.", 0, C_CHARCOAL),
    ], size=13)

    s = make_content_slide(prs, "7.4 Streaming run #1 (MEDIQ B=100 R=10 seed=0)", kicker="실험 결과")
    _add_image(s, ASSETS / "per_round_delta.png", 0.6, 1.4, w_in=8.0)
    _add_bullets(s, 8.8, 1.5, 4.1, 5.6, [
        "wall ≈ 9 h 45 min",
        "4 / 10 ACCEPTS",
        "2 / 10 reject",
        "4 / 10 INVALID",
        ("INVALID 원인: (3개) max_agents=6 binding, (1개) orphaned-prune.", 0, C_ACCENT),
        ("legacy v2가 3 도메인 × 3 iter에서 1 ACCEPT였던 것과 비교 — streaming-mode가 fire함.", 0, C_FOREST),
    ], size=12)

    s = make_content_slide(prs, "7.5 Streaming run #1 — DAG 진화 시퀀스 + token cost", kicker="실험 결과")
    _add_image(s, ASSETS / "dag_evolution_seq.png", 0.2, 1.4, w_in=12.9)
    _add_image(s, ASSETS / "token_cost_pareto.png", 0.6, 4.4, w_in=4.5)
    _add_bullets(s, 5.5, 4.4, 7.3, 2.7, [
        "Test (n=50): CoT 68 % / P-E 58 % / Evolved 62 %.",
        "Token / task: CoT 0.46k / P-E 0.79k / Evolved 2.90k (~6.3×).",
        ("토큰을 더 써도 정확도가 올라가지 않음 (이 split에서).", 0, C_ACCENT),
        ("Evolved가 P-E보다 +4 pp이지만 CoT보다 −6 pp.", 0, C_CHARCOAL),
        ("→ 'streaming 승리'로 인용 금지 — publishable 신호는 paired ACCEPT.", 0, C_FOREST),
    ], size=11)

    # ---- §7.6 Streaming run #2 (NEW) ---------------------------------
    s = make_content_slide(prs, "7.6 Streaming run #2 (MEDIQ s=1) — Δ vs P-E +4pp 재현",
                           kicker="실험 결과")
    _add_table(s, 0.6, 1.4, 12.1, 1.6, [
        ["", "run #1 (s=0)", "run #2 (s=1)"],
        ["ACCEPT / reject / INVALID", "4 / 2 / 4", "2 / 6 / 2"],
        ["wall", "9h 45m", "9h 32m"],
    ])
    _add_table(s, 0.6, 3.2, 12.1, 1.7, [
        ["", "run #1 (s=0)", "run #2 (s=1)"],
        ["CoT test (n=50)", "68 %", "46 %"],
        ["P-E test (n=50)", "58 %", "42 %"],
        ["Evolved test (n=50)", "62 %", "46 %"],
    ])
    _add_bullets(s, 0.6, 5.1, 12.1, 2.0, [
        ("절대 acc는 22 pp 변동 (CoT 68 → 46) — split 난이도가 다름.", 0, C_ACCENT),
        ("그러나 Δ vs P-E = +4pp 가 두 시드에서 *동일하게* 재현 — paired Δ가 split-난이도를 cancel.", 0, C_FOREST),
        ("INVALID 2건은 둘 다 orphan-edit 패턴 (run #1 4건의 동일 모드 재현) — v3에서 해결.", 0, C_CHARCOAL),
    ], size=13)

    # ---- §7.7 Patches + sanity (renumbered from §7.6) ----------------
    s = make_content_slide(prs, "7.7 §5.1.5 patches + sanity — streaming의 한계 인정",
                           kicker="실험 결과")
    _add_pill(s, 0.6, 1.4, 4.5, 0.5, "Patches (commit 8405c78)", fill=C_FOREST, size=12)
    _add_bullets(s, 0.6, 1.95, 6.0, 5.0, [
        "max_agents → user prompt # Constraints 라인",
        "AT CAP 시 add_agent 금지 명시",
        "prune-DAG reminder (orphan 방지)",
        "concept-level anti-repeat callout",
        "--max-agents default 6 → 8",
    ], size=12)
    _add_pill(s, 6.9, 1.4, 5.7, 0.5, "Sanity (sanity_streaming_v3_mediq_s0)", fill=C_FOREST, size=12)
    _add_bullets(s, 6.9, 1.95, 5.7, 5.0, [
        "B=20 R=3 mediq seed=0, ~50분 wall",
        "r2 paired ACCEPT (Δ=+5pp)",
        "_build_user_prompt 직접 호출로 # Constraints 렌더 확인",
        ("R=3 / cap=8에서 cap-binding / concept-level fire 미관찰 (sample 부족).", 0, C_ACCENT),
    ], size=12)
    _add_textbox(s, 0.6, 6.5, 12.1, 0.6,
                 "결론: streaming의 cross-batch 절대 acc 비교는 구조적으로 깨짐 → §4.4 v3로 이행.",
                 size=12, italic=True, color=C_ACCENT)

    # ---- §7.8 Controller v3 first run (NEW) --------------------------
    s = make_content_slide(prs, "7.8 Controller v3 첫 run — Evolved beats CoT (+6pp)",
                           kicker="실험 결과")
    _add_image(s, REPO_ROOT / "results" / "v3_mediq_s0" / "plots" / "accuracy_vs_iter.png",
               0.4, 1.45, w_in=6.0)
    _add_image(s, REPO_ROOT / "results" / "v3_mediq_s0" / "plots" / "arch_size.png",
               6.7, 1.45, w_in=6.0)
    _add_table(s, 0.6, 4.7, 12.1, 1.7, [
        ["Method", "Test acc (n=50)", "Tokens / task", "Δ vs CoT"],
        ["CoT", "46.0 %", "0.55k", "—"],
        ["P-E", "50.0 %", "1.09k", "+4 pp"],
        ["Evolved (v3, 8 ag, 15 ed)", "52.0 %", "3.28k (3.0× P-E)", "+6 pp ✓"],
    ])
    _add_textbox(s, 0.6, 6.45, 12.1, 0.4,
                 "MEDIQ seed=0, n_train=30, max_iters=10, wall ≈ 5h (v2 streaming의 ½)",
                 size=11, italic=True, color=C_CHARCOAL)
    _add_textbox(s, 0.6, 6.85, 12.1, 0.4,
                 "4 ACCEPT (iter 1, 5, 7, 8) / 6 reject / 0 INVALID — orphan auto-drop이 streaming 30-40% INVALID를 완전 해소.",
                 size=11, italic=True, color=C_FOREST)

    # ---- §7.9 v3 vs v2 streaming summary (NEW) -----------------------
    s = make_content_slide(prs, "7.9 v3 vs v2 streaming — 종합 (3 run, MEDIQ)",
                           kicker="실험 결과")
    _add_table(s, 0.6, 1.4, 12.1, 2.0, [
        ["Run", "wall", "ACCEPT", "INVALID", "Final n_ag", "Δ vs CoT", "Δ vs P-E"],
        ["v2 streaming run #1 (s=0)", "9h 45m", "4 / 10", "4 / 10", "6", "−6 pp", "+4 pp"],
        ["v2 streaming run #2 (s=1)", "9h 32m", "2 / 10", "2 / 10", "4", "0 pp", "+4 pp"],
        ["v3 first run (s=0)", "5h", "4 / 10", "0 / 10", "8", "+6 pp ✓", "+2 pp"],
    ])
    _add_bullets(s, 0.6, 3.7, 12.1, 3.4, [
        ("v3가 같은 ACCEPT 횟수를 ½ wall에 / INVALID 0 / 더 큰 그래프 (8 ag) 로 산출.", 0, C_FOREST),
        ("Δ vs CoT = +6pp 는 파일럿 전체에서 처음 — 단일 시드 한정 강신호.", 0, C_FOREST),
        ("Δ vs P-E 는 v3에서 +2pp로 줄었지만, 이는 P-E acc가 50%로 좋게 나온 split 효과 가능.", 0, C_CHARCOAL),
        ("Robust 신호: streaming 두 시드 paired Δ vs P-E = +4pp 재현 (run #1·#2).", 0, C_CHARCOAL),
        ("미해결: v3는 단일 시드. multi-seed × multi-domain 재현이 venue submission 직전 핵심 작업.", 0, C_ACCENT),
    ], size=13)

    # ---- Analysis & Discussion (8.1, 8.2) ------------------------------
    make_section_slide(prs, "8. 분석 (Analysis)",
                       "측정 노이즈 · 토큰 비용 · 실패 모드")

    s = make_content_slide(prs, "8.1 측정 노이즈 — paired-batch / full-pass 어떻게 다루나", kicker="분석")
    _add_bullets(s, 0.6, 1.4, 12.1, 5.6, [
        "Same-graph cross-run 노이즈 ±13 pp (FinanceBench v2 retry, MEDIQ run #1→#2 22pp 변동).",
        "Legacy mode는 cross-batch 비교 → 작은 architectural Δ가 노이즈에 휩쓸려 reject.",
        ("Streaming paired-batch는 best vs candidate를 SAME batch로 평가해 cancel — "
         "1/9 → 4/10 ACCEPT, 두 시드 paired Δ vs P-E = +4pp 재현.", 0, C_FOREST),
        ("그러나 streaming은 absolute acc가 bootstrap 분포에 의존 → cross-batch best_val_acc 비신뢰.", 0, C_ACCENT),
        ("v3 full-pass는 동일한 30 train task 위에서 strict 비교 가능 — MEDIQ s=0에서 "
         "INVALID 0 / 4 ACCEPT 산출.", 0, C_FOREST),
        ("Robust 신호 ranking: paired Δ vs P-E = +4pp (두 시드) > v3 +6pp vs CoT (단일 시드).", 0, C_CHARCOAL),
    ], size=13)

    s = make_content_slide(prs, "8.2 토큰 비용과 실패 모드 (3 모드 종합)", kicker="분석")
    _add_bullets(s, 0.6, 1.4, 12.1, 5.6, [
        ("Worker 토큰이 controller 토큰 대비 ~55× (streaming run #1: 4.03M : 73.7k).", 0, C_CHARCOAL),
        ("v3은 [QUERY] fire 시 노드당 1 → 3 LLM 호출 → worker 비용 평균 1.5-2× 증가, "
         "그러나 INVALID 0 으로 wall efficiency는 streaming의 ½.", 0, C_FOREST),
        ("실패 모드 (v2 run #1·#2 종합): cap binding (6 round), orphan prune (4 round), "
         "concept-level repeat (rounds 5-10).", 0, C_ACCENT),
        ("v3 해결: (1) max_agents prompt + apply_edits HARD plumbing, (2) orphan auto-drop, "
         "(3) concept-level repeat은 여전히 soft → orchestration-layer hard enforcement는 향후.", 0, C_FOREST),
        ("Cost-Pareto: v3 (52% @ 3.28k tokens) vs CoT (46% @ 0.55k); 6× 토큰에 +6pp — "
         "tokens-for-acc 트레이드오프는 도메인별로 검증 필요.", 0, C_CHARCOAL),
    ], size=12)

    # ---- Limitations & Future (10.1, 10.2) -----------------------------
    make_section_slide(prs, "9. 한계 및 향후 작업",
                       "정직한 자기 평가")

    s = make_content_slide(prs, "9.1 한계 (Limitations) — v3 반영", kicker="한계")
    _add_bullets(s, 0.6, 1.4, 12.1, 5.6, [
        ("v3 +6pp vs CoT 는 *단일 시드 단일 도메인* (MEDIQ s=0) — multi-seed × multi-domain 재현 미완.", 0, C_ACCENT),
        ("±13 pp 노이즈 floor 가 v3에서도 여전 (run-to-run absolute 22pp 변동 관찰됨).", 0, C_ACCENT),
        ("LLM-judge가 same-family Qwen → self-bias 위험 (FinanceBench, AgentClinic).", 0, C_ACCENT),
        ("단일 backbone (Qwen2.5-32B), 단일 judge.", 0, C_ACCENT),
        ("v3 wall 5h × seed × domain → 3×3 sweep ~45h (streaming 의 ½, 그러나 여전히 multi-day).", 0, C_ACCENT),
        ("ADAS / MaAS / Puppeteer 직접 비교 미실시 — reviewer 첫 질문이 될 항목.", 0, C_ACCENT),
        ("Concept-level anti-repeat이 v3에서도 prompt-only soft — orch-layer hard enforcement 필요.", 0, C_ACCENT),
        ("[QUERY] fire-rate / parse-rate / 효과 측정 미완 — v3 component-level ablation의 input.", 0, C_ACCENT),
    ], size=13)

    s = make_content_slide(prs, "9.2 향후 작업 (Future Work)", kicker="향후 작업")
    _add_bullets(s, 0.6, 1.4, 12.1, 5.6, [
        ("§5.2 v3 multi-seed sweep — MEDIQ s={1,2}, AgentClinic s={0,1,2}, "
         "FinanceBench s={0,1,2} (예상 ~45h, venue submission 직전 핵심).", 0, C_FOREST),
        "§5.3 random-persona ablation (v3 persona를 random text로 교체; reviewer 질문 #1).",
        "§5.4 harness ablation — controller {none, random, fixed-topo, full v3}.",
        "§5.5 [QUERY] / [SUMMARY] component ablation — 각 채널 끄고 효과 격리.",
        "§5.6 두 번째 backbone (Qwen3-72B + Claude/GPT API).",
        "§5.7 LLM-judge 교체 (Claude Haiku 4.5 또는 GPT-4.1-mini).",
        ("§5.8 ADAS + (MaAS / Puppeteer / EvoMAC) 직접 baseline — 리뷰어 필수 항목, 비협상.", 0, C_ACCENT),
    ], size=13)

    # ---- Conclusion ----------------------------------------------------
    s = make_content_slide(prs, "10. 결론 (Conclusion)", kicker="결론")
    _add_bullets(s, 0.6, 1.4, 12.1, 5.6, [
        ("Reflection-only 멀티-에이전트 진화는 *그래프를 움직이게* 만든다 — v2가 도메인 어휘를 가진 "
         "specialist persona / triage department 구조를 출력함이 그 증거.", 0, C_FOREST),
        ("v2 @ n=30 단일 시드는 baseline을 못 이김 — ±13pp 노이즈가 architectural 효과를 가림.", 0, C_ACCENT),
        ("Streaming paired-batch가 노이즈 floor를 부분 해결 — 1/9 → 4/10 ACCEPT, "
         "두 시드 paired Δ vs P-E = +4pp 재현.", 0, C_FOREST),
        ("Controller v3 (full-pass + 샘플 단위 reflection 30→3→1 + side-channel Q&A + "
         "orphan auto-drop) 는 INVALID 0, wall ½, MEDIQ s=0에서 *처음으로* CoT 를 +6pp로 이김.", 0, C_FOREST),
        ("미해결: 단일 시드 한정 강신호. multi-seed × multi-domain × 두 번째 백본 × "
         "ADAS 직접 비교가 venue submission 직전 잔여 작업.", 0, C_CHARCOAL),
    ], size=13)

    # ---- Q&A ---------------------------------------------------------
    make_conclusion_slide(
        prs,
        "감사합니다 — 질의응답",
        "코드 + 결과: github.com/namam3gy · 로드맵: references/roadmap_ko.md",
    )

    OUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT_PATH))
    print(f"wrote {OUT_PATH} ({len(prs.slides)} slides)")


if __name__ == "__main__":
    build()
