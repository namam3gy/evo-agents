"""Build PPT 1 — Advisor talk (English, ~34 slides), Forest & Moss palette.

Uses python-pptx. Outputs to docs/presentations/2026-04-26-evo-agents/1-advisor-talk.pptx.
Run via:
    uv run python scripts/build_advisor_deck.py
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

REPO_ROOT = Path(__file__).resolve().parents[1]
DECK_DIR = REPO_ROOT / "docs" / "presentations" / "2026-04-26-evo-agents"
ASSETS = DECK_DIR / "assets"
OUT_PATH = DECK_DIR / "1-advisor-talk.pptx"

# ---------------------------------------------------------------------------
# palette (Forest & Moss)
# ---------------------------------------------------------------------------
C_FOREST = RGBColor(0x2C, 0x5F, 0x2D)
C_MOSS = RGBColor(0x97, 0xBC, 0x62)
C_CREAM = RGBColor(0xF5, 0xF5, 0xF5)
C_CHARCOAL = RGBColor(0x36, 0x45, 0x4F)
C_ACCENT = RGBColor(0xB8, 0x50, 0x42)
C_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
C_LIGHT_GREY = RGBColor(0xCF, 0xD3, 0xD7)
C_NEAR_BLACK = RGBColor(0x1A, 0x22, 0x2A)

FONT_HEADER = "Calibri"
FONT_BODY = "Calibri Light"

SLIDE_W_IN = 13.333
SLIDE_H_IN = 7.5

# ---------------------------------------------------------------------------
# helpers
# ---------------------------------------------------------------------------


def _set_slide_background(slide, color: RGBColor) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_textbox(slide, x_in, y_in, w_in, h_in, text,
                 *, size=14, color=C_CHARCOAL, bold=False, italic=False,
                 font=FONT_BODY, align=PP_ALIGN.LEFT,
                 anchor=MSO_ANCHOR.TOP, line_spacing=1.15):
    tx = slide.shapes.add_textbox(Inches(x_in), Inches(y_in),
                                  Inches(w_in), Inches(h_in))
    tf = tx.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Inches(0.06)
    tf.margin_right = Inches(0.06)
    tf.margin_top = Inches(0.04)
    tf.margin_bottom = Inches(0.04)
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = line_spacing
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    r.font.name = font
    return tx


def _add_bullets(slide, x_in, y_in, w_in, h_in, items,
                 *, size=15, color=C_CHARCOAL, bullet="•",
                 font=FONT_BODY, line_spacing=1.25, sub_size=13):
    """Add a bulleted list. Items can be (text, indent_level, optional_color)."""
    tx = slide.shapes.add_textbox(Inches(x_in), Inches(y_in),
                                  Inches(w_in), Inches(h_in))
    tf = tx.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.06)
    tf.margin_right = Inches(0.06)
    tf.margin_top = Inches(0.04)
    tf.margin_bottom = Inches(0.04)

    for i, raw in enumerate(items):
        if isinstance(raw, str):
            text, lvl, run_color = raw, 0, color
        else:
            text = raw[0]
            lvl = raw[1] if len(raw) > 1 else 0
            run_color = raw[2] if len(raw) > 2 else color
        prefix = "  " * lvl + (bullet if lvl == 0 else "–") + " "
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = line_spacing
        r = p.add_run()
        r.text = prefix + text
        r.font.size = Pt(size if lvl == 0 else sub_size)
        r.font.color.rgb = run_color
        r.font.name = font
    return tx


def _add_title(slide, title_text, *, color=C_CHARCOAL, size=32, x_in=0.6, y_in=0.4,
               w_in=12.1, h_in=0.9):
    return _add_textbox(slide, x_in, y_in, w_in, h_in, title_text,
                        size=size, color=color, bold=True, font=FONT_HEADER,
                        align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)


def _add_accent_strip(slide, color=C_FOREST, height_in=0.18):
    """Thick top strip used as visual motif on content slides."""
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0,
                                  Inches(SLIDE_W_IN), Inches(height_in))
    rect.fill.solid()
    rect.fill.fore_color.rgb = color
    rect.line.fill.background()


def _add_image(slide, path: Path, x_in, y_in, w_in, h_in=None):
    if h_in is None:
        return slide.shapes.add_picture(str(path), Inches(x_in), Inches(y_in),
                                        width=Inches(w_in))
    return slide.shapes.add_picture(str(path), Inches(x_in), Inches(y_in),
                                    width=Inches(w_in), height=Inches(h_in))


def _add_pill(slide, x_in, y_in, w_in, h_in, text,
              *, fill=C_MOSS, font_color=C_WHITE, size=12, bold=True):
    rect = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  Inches(x_in), Inches(y_in),
                                  Inches(w_in), Inches(h_in))
    rect.fill.solid()
    rect.fill.fore_color.rgb = fill
    rect.line.fill.background()
    tf = rect.text_frame
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.03)
    tf.margin_bottom = Inches(0.03)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = font_color
    r.font.name = FONT_HEADER
    return rect


def _add_table(slide, x_in, y_in, w_in, h_in, data, *,
               header_fill=C_FOREST, header_font=C_WHITE,
               body_font=C_CHARCOAL, body_size=12,
               header_size=13, alt_fill=C_CREAM):
    rows = len(data)
    cols = len(data[0])
    table = slide.shapes.add_table(rows, cols,
                                   Inches(x_in), Inches(y_in),
                                   Inches(w_in), Inches(h_in)).table
    for j, cell_text in enumerate(data[0]):
        cell = table.cell(0, j)
        cell.fill.solid()
        cell.fill.fore_color.rgb = header_fill
        tf = cell.text_frame
        tf.margin_left = Inches(0.05)
        tf.margin_right = Inches(0.05)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = str(cell_text)
        r.font.size = Pt(header_size)
        r.font.bold = True
        r.font.color.rgb = header_font
        r.font.name = FONT_HEADER
    for i in range(1, rows):
        row_alt = alt_fill if i % 2 else C_WHITE
        for j, cell_text in enumerate(data[i]):
            cell = table.cell(i, j)
            cell.fill.solid()
            cell.fill.fore_color.rgb = row_alt
            tf = cell.text_frame
            tf.margin_left = Inches(0.06)
            tf.margin_right = Inches(0.06)
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER if j > 0 else PP_ALIGN.LEFT
            r = p.add_run()
            r.text = str(cell_text)
            r.font.size = Pt(body_size)
            r.font.bold = False
            r.font.color.rgb = body_font
            r.font.name = FONT_BODY
    return table


# ---------------------------------------------------------------------------
# slide layouts
# ---------------------------------------------------------------------------


def make_title_slide(prs, title, subtitle, footer):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_background(slide, C_FOREST)
    # decorative moss bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                 Inches(0), Inches(3.0),
                                 Inches(0.3), Inches(1.5))
    bar.fill.solid()
    bar.fill.fore_color.rgb = C_MOSS
    bar.line.fill.background()
    _add_textbox(slide, 0.9, 2.5, 12.0, 1.4, title,
                 size=44, color=C_WHITE, bold=True, font=FONT_HEADER,
                 anchor=MSO_ANCHOR.MIDDLE)
    _add_textbox(slide, 0.9, 4.0, 12.0, 0.8, subtitle,
                 size=20, color=C_LIGHT_GREY, font=FONT_BODY, italic=True,
                 anchor=MSO_ANCHOR.TOP)
    _add_textbox(slide, 0.9, 6.7, 12.0, 0.5, footer,
                 size=12, color=C_LIGHT_GREY, font=FONT_BODY,
                 anchor=MSO_ANCHOR.TOP)
    return slide


def make_section_slide(prs, label, body_text=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_background(slide, C_FOREST)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                 Inches(0.6), Inches(3.4),
                                 Inches(0.25), Inches(0.9))
    bar.fill.solid()
    bar.fill.fore_color.rgb = C_MOSS
    bar.line.fill.background()
    _add_textbox(slide, 1.1, 3.3, 12.0, 1.0, label,
                 size=42, color=C_WHITE, bold=True, font=FONT_HEADER,
                 anchor=MSO_ANCHOR.MIDDLE)
    if body_text:
        _add_textbox(slide, 1.1, 4.4, 12.0, 1.4, body_text,
                     size=18, color=C_LIGHT_GREY, font=FONT_BODY, italic=True,
                     anchor=MSO_ANCHOR.TOP)
    return slide


def make_content_slide(prs, title, *, kicker=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_background(slide, C_CREAM)
    _add_accent_strip(slide, color=C_FOREST)
    if kicker:
        _add_textbox(slide, 0.6, 0.27, 12.0, 0.32, kicker,
                     size=11, color=C_MOSS, bold=True, font=FONT_HEADER,
                     anchor=MSO_ANCHOR.TOP)
    _add_title(slide, title, y_in=0.55)
    return slide


def make_conclusion_slide(prs, big, sub):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_background(slide, C_FOREST)
    _add_textbox(slide, 0.6, 2.6, 12.1, 1.3, big,
                 size=54, color=C_WHITE, bold=True, font=FONT_HEADER,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    _add_textbox(slide, 0.6, 4.1, 12.1, 1.6, sub,
                 size=20, color=C_LIGHT_GREY, font=FONT_BODY,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP)
    return slide


# ---------------------------------------------------------------------------
# concrete slides
# ---------------------------------------------------------------------------


def build():
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W_IN)
    prs.slide_height = Inches(SLIDE_H_IN)

    # --- 1. Title --------------------------------------------------------
    make_title_slide(
        prs,
        "Self-Evolving Multi-Agent DAGs from Pure Reflection",
        "A pilot study on reflection-only graph evolution — no search, no RL",
        "thyun.park · advisor talk · 2026-04-26",
    )

    # --- 2. TL;DR --------------------------------------------------------
    s = make_content_slide(prs, "TL;DR — what we tested",
                           kicker="One-liner")
    _add_textbox(s, 0.6, 1.5, 12.1, 1.0,
                 "Can a multi-agent DAG (topology + personas + edges) be progressively edited "
                 "by an LLM controller using ONLY in-context reflection over trajectory tapes — "
                 "no search archive, no MCTS, no RL?",
                 size=18, color=C_CHARCOAL, italic=True, line_spacing=1.3)
    _add_bullets(s, 0.6, 2.9, 12.1, 3.7, [
        "Backbone: Qwen2.5-32B-Instruct (vLLM, single H200, bf16)",
        "Domains: FinanceBench / MEDIQ / AgentClinic (post-pivot from GSM8K)",
        "Methods compared: CoT · Planner-Executor · Evolved DAG (controller v2)",
        "Scoring: test accuracy + paired-accept rate (under streaming evolve)",
        ("Controller v1 (n=30): at-or-below baselines on test, 'verifier-add' reflex (H1 weak fail)", 1, C_ACCENT),
        ("Controller v2 (n=30): qualitatively new behavior; test win not yet (noise dominates)", 1, C_ACCENT),
        ("Streaming evolve mode: 4 / 10 paired ACCEPTS in run #1 (vs 1 / 9 legacy)", 1, C_FOREST),
    ], size=15)
    _add_pill(s, 0.6, 6.7, 4.5, 0.45, "§5.2 multi-seed sweep — running NOW",
              fill=C_ACCENT, font_color=C_WHITE, size=12)

    # --- 3. Section: Setup ----------------------------------------------
    make_section_slide(prs, "Setup", "DAG formalism · baselines · pilot infrastructure")

    # --- 4. DAG formalism + baselines (figure) ---------------------------
    s = make_content_slide(prs, "DAG formalism + the three graph shapes")
    _add_image(s, ASSETS / "dag_baselines.png", 0.5, 1.5, w_in=12.3)
    _add_bullets(s, 0.6, 5.2, 12.1, 1.9, [
        "Graph G = (Agents A, Edges E ⊆ A × A) with reserved START / END nodes; orchestrator runs in topological order.",
        "An agent = (name, persona text). The persona is what the LLM is conditioned on at this node.",
        "Three shapes co-exist: CoT (1 agent), Planner-Executor seed (2 agents, baseline), Evolved (controller-edited).",
    ], size=14)

    # --- 5. Pilot infrastructure -----------------------------------------
    s = make_content_slide(prs, "Pilot infrastructure (single repo, runnable end-to-end)",
                           kicker="Implementation")
    _add_bullets(s, 0.6, 1.5, 6.0, 5.6, [
        "src/llm.py · OpenAI-compatible chat client (vLLM at :8000)",
        "src/graph.py · Graph + edit ops (add/remove agent + edge, rewrite_persona)",
        "src/orchestrator.py · runs the graph on a Task → Tape",
        "src/controller.py · Reflection LLM that emits an EditBatch",
        "src/evolve.py · legacy + streaming evolve loops",
        "src/score.py · MCQ exact-match (MEDIQ) + LLM-judge (FB / AC)",
        "scripts/run_pilot.py · driver (baselines + evolve + test eval)",
        "uv-managed env, vLLM 0.19.1 patches, gcc auto-install in serve_vllm.sh",
    ], size=14)
    _add_pill(s, 7.0, 1.5, 5.7, 0.5, "Backbone: Qwen2.5-32B-Instruct (bf16, single H200)",
              fill=C_FOREST, size=13)
    _add_bullets(s, 7.0, 2.2, 5.7, 4.9, [
        "vLLM with prefix caching",
        "EVO_GPU_UTIL=0.55 (shared GPU; ~78 GB)",
        "max_model_len = 16384 (controller prompts + tapes)",
        "EVO_MODEL / EVO_BASE_URL env-overridable",
        "Each run writes results/<run_id>/{evolve_log, results, plots/}",
        "Standard plots: accuracy_vs_iter / arch_size / edit_mix",
    ], size=13)

    # --- 6. Section: Calibration -----------------------------------------
    make_section_slide(prs, "Calibration on GSM8K (calib_01)",
                       "First end-to-end run · what it taught us · why we pivoted")

    # --- 7. calib_01 setup + result --------------------------------------
    s = make_content_slide(prs, "calib_01 — first calibrated run on GSM8K",
                           kicker="Stage 1 — Calibration")
    _add_image(s, REPO_ROOT / "results" / "calib_01" / "plots" / "accuracy_vs_iter.png",
               0.6, 1.5, w_in=6.4)
    _add_bullets(s, 7.3, 1.5, 5.6, 5.6, [
        "Setup: n_val = n_test = 50, seed = 0, max_iters = 3, B = ~66 min wall.",
        "Headline: Evolved 86 % test, CoT 92 %, P-E 90 %.",
        ("Evolved underperforms BOTH baselines on test, at 2–3.7× more tokens.", 1, C_ACCENT),
        "Val saturated (CoT = P-E = 94 %) → no discriminating signal at n = 50.",
        "Controller did emit hypothetical edits — graph mutates, but not towards a useful target.",
        "Pipeline finding: under old accept policy, best_graph and best_val_acc could disagree.",
    ], size=13)

    # --- 8. Insight 1 - Opt-2 strict --------------------------------------
    s = make_content_slide(prs, "Insight 1 — Opt-2 strict accept rule",
                           kicker="Stage 1 — Calibration")
    _add_bullets(s, 0.6, 1.4, 12.1, 5.6, [
        ("Problem: under accept_slack > 0, best_graph could be replaced by a near-best graph "
         "while best_val_acc stayed at the original peak — the two fields disagreed.", 0, C_ACCENT),
        "Fix (Opt-2 strict): both best_graph and best_val_acc advance only when val_acc > best_val_acc.",
        "Rejected candidates do NOT replace best_graph; ties also REJECT (default accept_slack = 0.0).",
        "Verified on results/smoke_opt2/: tied evolved iterations REJECTED; final best_graph == seed.",
        ("Why it matters for the rest of the pilot: we need best_graph to be the graph that actually "
         "achieved the val score we report.", 0, C_FOREST),
    ], size=14)

    # --- 9. Insight 2 — GSM8K wrong testbed ------------------------------
    s = make_content_slide(prs, "Insight 2 — GSM8K is the wrong testbed",
                           kicker="Stage 1 — Calibration")
    _add_bullets(s, 0.6, 1.4, 12.1, 5.6, [
        "Single-model saturation: Qwen2.5-32B already hits ~94 % CoT — no headroom for multi-agent gain.",
        "Self-contained text: every problem is fully specified — no information asymmetry between agents.",
        "Linear arithmetic: planner / executor decomposition adds no leverage over a single CoT pass.",
        ("→ Multi-agent edits cannot exploit a structural advantage that does not exist in this benchmark.", 0, C_ACCENT),
        ("Decision (2026-04-24): pivot to information-asymmetric domains where multi-agent has a real lever.", 0, C_FOREST),
    ], size=14)

    # --- 10. Section: domain pivot ---------------------------------------
    make_section_slide(prs, "Domain pivot — three new benchmarks",
                       "FinanceBench · MEDIQ · AgentClinic")

    # --- 11. Three benchmarks ---------------------------------------------
    s = make_content_slide(prs, "FinanceBench / MEDIQ / AgentClinic",
                           kicker="Stage 2 — Domain pivot")
    _add_pill(s, 0.6, 1.5, 4.0, 0.5, "FinanceBench (LLM-judge)", fill=C_FOREST, size=13)
    _add_bullets(s, 0.6, 2.05, 4.0, 5.0, [
        "10-K extraction QA",
        "GAAP / TTM / fiscal year",
        "Long-context retrieval",
        "Failure: hedging, period mix-up",
    ], size=12)
    _add_pill(s, 4.7, 1.5, 4.0, 0.5, "MEDIQ (MCQ exact-match)", fill=C_FOREST, size=13)
    _add_bullets(s, 4.7, 2.05, 4.0, 5.0, [
        "Clinical vignette → MCQ",
        "Differential diagnosis",
        "Base-rate reasoning",
        "Failure: demographic blindness",
    ], size=12)
    _add_pill(s, 8.8, 1.5, 4.0, 0.5, "AgentClinic (LLM-judge)", fill=C_FOREST, size=13)
    _add_bullets(s, 8.8, 2.05, 4.0, 5.0, [
        "Single-pass clinical case",
        "Specialty routing",
        "Decisive answer",
        "Failure: triage error, hedging",
    ], size=12)
    _add_textbox(s, 0.6, 6.5, 12.1, 0.6,
                 "Sanity n=3 per benchmark passed (results/sanity_{financebench,mediq,agentclinic}/) — "
                 "pipeline runs end-to-end on all three.",
                 size=12, italic=True, color=C_MOSS, font=FONT_BODY)

    # --- 12. Section v1 controller --------------------------------------
    make_section_slide(prs, "v1 controller @ n = 30",
                       "First domain-pivot measurement")

    # --- 13. v1 n=30 table ----------------------------------------------
    s = make_content_slide(prs, "v1 controller @ n = 30 — three-domain measurement",
                           kicker="Stage 3 — v1 controller")
    _add_table(s, 0.6, 1.6, 12.1, 2.5, [
        ["Domain", "CoT test", "P-E test", "Evolved test", "Δ vs best baseline"],
        ["FinanceBench", "70.0 %", "66.7 %", "70.0 %", "0 pp"],
        ["MEDIQ", "43.3 %", "43.3 %", "50.0 %", "+6.7 pp (within ±18 pp noise)"],
        ["AgentClinic", "66.7 %", "70.0 %", "70.0 %", "0 pp"],
    ])
    _add_bullets(s, 0.6, 4.4, 12.1, 2.6, [
        "Setup: --n-train 10 --n-val 30 --n-test 30 --max-iters 3 --seed 0.",
        "FinanceBench: same `add_verifier` edit emitted 3 rounds in a row — no domain vocabulary.",
        "MEDIQ: edits varied after iter 2 was ACCEPTED; some specialty mention.",
        "AgentClinic: alternated `summarizer ↔ verifier`.",
        ("Read: H1 weakly falsified at this controller version. Headroom likely exists, "
         "but v1 is too thin to exploit it.", 0, C_ACCENT),
    ], size=13)

    # --- 14. Insight: verifier reflex -----------------------------------
    s = make_content_slide(prs, "Insight — the v1 'verifier-add' reflex",
                           kicker="Stage 3 — v1 controller")
    _add_bullets(s, 0.6, 1.4, 12.1, 5.6, [
        "Across all three domains, v1's most common edit was a generic `add_agent('verifier')`.",
        "Rationales were plausible-sounding but content-free — no GAAP, no differential diagnosis, no triage.",
        ("v1 system prompt was only ~30 lines: 'edit the graph, here are the ops'. "
         "No notion of domain expertise, no persona-authoring rules, no anti-repeat.", 0, C_CHARCOAL),
        ("Hypothesis (drives H2): a thicker, domain-grounded controller can produce specialist personas "
         "that the v1 controller architecturally cannot.", 0, C_FOREST),
    ], size=14)

    # --- 15. Section v2 ---------------------------------------------------
    make_section_slide(prs, "Controller v2 — organization designer",
                       "System-prompt redesign · domain briefs · specialist personas")

    # --- 16. v1 vs v2 loop figure ---------------------------------------
    s = make_content_slide(prs, "v1 vs v2 — same loop, different conditioning",
                           kicker="Stage 4 — v2 controller")
    _add_image(s, ASSETS / "controller_v1_v2_loop.png", 0.6, 1.5, w_in=12.1)
    _add_textbox(s, 0.6, 6.5, 12.1, 0.6,
                 "Same reflection loop — only the SYSTEM prompt and the side-channel domain brief differ. "
                 "Output style flips from generic 'add verifier' to cited specialist personas.",
                 size=12, italic=True, color=C_MOSS)

    # --- 17. What's new in v2 ---------------------------------------------
    s = make_content_slide(prs, "What's new in v2",
                           kicker="Stage 4 — v2 controller")
    _add_bullets(s, 0.6, 1.4, 12.1, 5.6, [
        "Org-designer framing: 'design an organization of domain experts that solves these tasks better than the seed P-E.'",
        "Mandatory persona-authoring rules: each new agent must cite specialty + concrete procedure; generic verifier/summarizer FORBIDDEN unless paired with specialty.",
        "Domain briefs (data/briefs/{financebench,mediq,agentclinic}.md, ~80–110 lines): task style, failure modes, useful expertise, anti-patterns.",
        "Active prune incentive: encourage `remove_agent` when an agent contributes nothing.",
        "Anti-repeat rule (string-level): 'do not repeat a rejected edit' (will be tightened to concept-level later).",
        "v2 sanity n=10 across 3 domains: specialist names appear (gaap_analyst, period_validator, differential_diagnostician, adolescent_specialist, …) — qualitative shift confirmed.",
    ], size=13)

    # --- 18. v2 n=30 table ------------------------------------------------
    s = make_content_slide(prs, "v2 @ n = 30 — measurement noise dominates",
                           kicker="Stage 4 — v2 controller")
    _add_table(s, 0.6, 1.6, 12.1, 2.5, [
        ["Domain", "CoT test", "P-E test", "Evolved test", "Δ vs best baseline", "best_graph"],
        ["FinanceBench (16k)", "73.3 %", "70.0 %", "83.3 %*", "+10 pp*", "seed (all REJECT)"],
        ["MEDIQ", "43.3 %", "46.7 %", "43.3 %", "−3.4 pp", "3-agent (iter 2 ACCEPT)"],
        ["AgentClinic", "60.0 %", "73.3 %", "66.7 %", "−6.6 pp", "seed (all REJECT)"],
    ])
    _add_bullets(s, 0.6, 4.4, 12.1, 2.7, [
        ("FinanceBench's +10 pp is SAME-GRAPH noise: best_graph == seed, yet test acc moved 13 pp from "
         "the simultaneous P-E run. Cannot be cited as a v2 win.", 0, C_ACCENT),
        ("AgentClinic iter-3 (REJECTED but notable): proposed triage → {gastro, cardio} → END. A literal "
         "specialty-department department; rejected on val tie under Opt-2 strict.", 0, C_CHARCOAL),
        ("Read: H2 BEHAVIOR satisfied (specialist personas, varied edits, prune); H2 TEST WIN not yet.", 0, C_FOREST),
    ], size=13)

    # --- 19. Insight: same-graph noise -----------------------------------
    s = make_content_slide(prs, "Insight — same-graph cross-run noise is ±13 pp",
                           kicker="Stage 4 — v2 controller")
    _add_bullets(s, 0.6, 1.4, 12.1, 5.6, [
        "FinanceBench v2 evolved test = 83.3 %, but best_graph was the seed P-E.",
        "Yet a parallel P-E test run on the same split scored 70 %. Same graph, two different test draws.",
        ("Difference 13 pp comes from: (a) vLLM batch-ordering, (b) KV-cache state non-determinism, "
         "(c) tqdm-driven scheduling differences, (d) sample selection at n = 30.", 0, C_CHARCOAL),
        ("→ Architectural effect on test ≤ ±13 pp at n = 30 cannot be distinguished from re-run noise. "
         "The signal we want is BELOW the noise floor at this measurement scale.", 0, C_ACCENT),
        ("Practical consequence: under Opt-2 strict + per-iter wall budget, most architectural changes "
         "get rejected before they can be evaluated over multiple noise-averaged batches. Drives the streaming pivot.", 0, C_FOREST),
    ], size=13)

    # --- 20. Section: Streaming -------------------------------------------
    make_section_slide(prs, "Streaming evolve mode",
                       "Bootstrap-sampled mini-batches · paired comparison · run #1 on MEDIQ")

    # --- 21. Why streaming + paired flow ---------------------------------
    s = make_content_slide(prs, "Why streaming — paired-batch comparison kills cross-batch noise",
                           kicker="Stage 5 — Streaming evolve")
    _add_image(s, ASSETS / "paired_batch_flow.png", 0.4, 1.5, w_in=12.5)
    _add_bullets(s, 0.6, 5.5, 12.1, 1.7, [
        "Each round draws a bootstrap mini-batch (B tasks) from train + val pool. Both best and candidate are evaluated on the SAME batch.",
        "Decision: ACCEPT iff c_acc > b_acc + ε. Δ is INSIDE-batch, so cross-batch difficulty drift cancels.",
        "Wall scales linearly with B: at B=100, ~52 min/round; at B=20, ~10 min/round.",
    ], size=13)

    # --- 22. Run #1 paired Δ ---------------------------------------------
    s = make_content_slide(prs, "Run #1 paired Δ — streaming fires",
                           kicker="Stage 5 — Streaming evolve")
    _add_image(s, ASSETS / "per_round_delta.png", 0.6, 1.4, w_in=8.4)
    _add_bullets(s, 9.2, 1.5, 3.7, 5.6, [
        "MEDIQ B=100 R=10 seed=0",
        "wall ≈ 9 h 45 min",
        "4 / 10 ACCEPTS",
        "2 / 10 reject",
        "4 / 10 INVALID",
        ("INVALIDs all 'max_agents reached' at cap=6 (3 of 4) or orphaned-prune (1 of 4).",
         0, C_ACCENT),
        ("Compare: v2 legacy across 3 domains × 3 iters = 1 ACCEPT total.", 0, C_FOREST),
    ], size=13)

    # --- 23. DAG evolution sequence --------------------------------------
    s = make_content_slide(prs, "Run #1 DAG evolution — five snapshots",
                           kicker="Stage 5 — Streaming evolve")
    _add_image(s, ASSETS / "dag_evolution_seq.png", 0.2, 1.7, w_in=12.9)
    _add_textbox(s, 0.6, 6.0, 12.1, 1.1,
                 "Specialist personas accumulate through R4 (differential_diagnostician → epidemiology_consultant → "
                 "adolescent_specialist). R5–R7 are reject / INVALID. R8 ACCEPTS a planner-prune + differential_generator — "
                 "an organisational reorg, not just an addition.",
                 size=12, italic=True, color=C_CHARCOAL)

    # --- 24. Token vs accuracy + test ------------------------------------
    s = make_content_slide(prs, "Token cost vs test accuracy — streaming run #1",
                           kicker="Stage 5 — Streaming evolve")
    _add_image(s, ASSETS / "token_cost_pareto.png", 0.6, 1.4, w_in=7.8)
    _add_bullets(s, 8.6, 1.5, 4.3, 5.6, [
        "CoT 68 % @ 0.46k tok",
        "P-E 58 % @ 0.79k tok",
        "Evolved 62 % @ 2.90k tok",
        ("More tokens did NOT buy accuracy on this MEDIQ split.", 0, C_ACCENT),
        ("Evolved beats P-E by +4 pp at 6.3× tokens.", 1, C_CHARCOAL),
        ("Evolved loses to CoT by 6 pp at 6.3× tokens.", 1, C_CHARCOAL),
        ("Don't cite this number as a 'streaming wins' result. The publishable signal is paired ACCEPT activity.",
         0, C_FOREST),
    ], size=12)

    # --- 25. Three blockers ----------------------------------------------
    s = make_content_slide(prs, "Three blockers identified before §5.2",
                           kicker="Stage 5 — Streaming evolve")
    _add_pill(s, 0.6, 1.5, 4.0, 0.5, "(1) Pass criterion broken", fill=C_ACCENT, size=12)
    _add_bullets(s, 0.6, 2.05, 4.0, 5.0, [
        ("`best_val_acc > seed_batch_acc` compares MAX over independent bootstrap batches "
         "of varying difficulty — not a quality signal.", 0, C_CHARCOAL),
        "Fix: drop it. Score on test acc + paired-accept rate.",
    ], size=11)
    _add_pill(s, 4.7, 1.5, 4.0, 0.5, "(2) max_agents=6 binds", fill=C_ACCENT, size=12)
    _add_bullets(s, 4.7, 2.05, 4.0, 5.0, [
        "After R4 the graph has 6 agents; controller (which doesn't see the cap) keeps proposing add_agent → INVALID.",
        "3 of 4 INVALID rounds are this exact pattern (~30 min/round burned).",
        "Fix: surface the cap in the prompt + raise default to 8.",
    ], size=11)
    _add_pill(s, 8.8, 1.5, 4.0, 0.5, "(3) Anti-repeat is string-level", fill=C_ACCENT, size=12)
    _add_bullets(s, 8.8, 2.05, 4.0, 5.0, [
        "Controller renamed the same 'case feature extractor' role 5×: differential_generator → clinical_filter → pediatrician.",
        "Anti-repeat rule reads: 'don't repeat a REJECTED edit' — which the model interprets at the string level.",
        "Fix: add concept-level callout in the system prompt.",
    ], size=11)

    # --- 26. Section: Patches --------------------------------------------
    make_section_slide(prs, "§5.1.5 patches + sanity",
                       "Mechanical fixes for the three blockers · sanity validated")

    # --- 27. Patches detail ----------------------------------------------
    s = make_content_slide(prs, "Patches landed (commit 8405c78)",
                           kicker="Stage 6 — Pre-§5.2 patches")
    _add_bullets(s, 0.6, 1.4, 12.1, 5.6, [
        "src/controller.py: `propose_edits` takes `max_agents` kwarg.",
        ("`_build_user_prompt` opens with a `# Constraints` block "
         "(`max_agents = N, current n_agents = n. K agent slots remaining before the cap.`).", 1),
        ("When `n_agents == max_agents`, slack line becomes `AT CAP — only remove_agent / "
         "rewrite_persona / topology edits are allowed.`", 1),
        "SYSTEM prompt adds a prune-DAG reminder ('removing X must not orphan any agent that depended on X').",
        "SYSTEM prompt adds a concept-level anti-repeat callout ('rename ≠ different idea').",
        "src/evolve.py: legacy + streaming loops both thread `max_agents` through to `propose_edits`.",
        "scripts/run_pilot.py: --max-agents default 6 → 8 (fits triage + 2-3 specialists + answer).",
        ("Doc patch (4): `best_val_acc > seed_batch_acc` already replaced by test acc + paired-accept rate "
         "in §5.2 / pilot.md §8.4 — no further edits needed.", 0, C_FOREST),
    ], size=13)

    # --- 28. sanity_v3 results -------------------------------------------
    s = make_content_slide(prs, "sanity_streaming_v3_mediq_s0 — patches green for §5.2",
                           kicker="Stage 6 — Pre-§5.2 patches")
    _add_table(s, 0.6, 1.5, 12.1, 2.4, [
        ["Round", "b_acc", "c_acc", "Δpp", "Verdict", "Edits (brief)"],
        ["0", "60.0 %", "—", "—", "seed", "(planner+executor, 2 ag 4 ed)"],
        ["1", "60.0 %", "60.0 %", "0.0", "reject", "+differential_diagnostician (3 ag)"],
        ["2", "50.0 %", "55.0 %", "+5.0", "ACCEPT", "+differential_diagnostician (same as r1)"],
        ["3", "65.0 %", "65.0 %", "0.0", "reject", "+base_rate_consultant (would be 4 ag)"],
    ])
    _add_bullets(s, 0.6, 4.2, 12.1, 2.9, [
        "Wall ≈ 50 min total (B=20 R=3, baselines + n_test=30 eval). Per-round wall ~10 min — linear scale from B=100's 52 min.",
        "Direct `_build_user_prompt(max_agents=8, n_agents=2)` invocation confirms the `# Constraints` block renders.",
        "Test n=30: CoT 43.3 % / P-E 56.7 % / Evolved 43.3 %. Sanity is for *patches function*, not test win.",
        ("Cap-binding (n_ag never reaches 8) and concept-level anti-repeat firing not exercised at R=3 — defer to §5.2.", 0, C_ACCENT),
        ("Round 2 ACCEPTED an exact-string-repeat of round 1's rejected edit on a different bootstrap batch — "
         "the soft-constraint anti-repeat did not inhibit. Will be re-examined at B=100 R=10.", 0, C_CHARCOAL),
    ], size=12)

    # --- 29. §5.2 ongoing ------------------------------------------------
    s = make_content_slide(prs, "§5.2 — ongoing right now (MEDIQ seed=1)",
                           kicker="Stage 7 — Live experiment")
    _add_pill(s, 0.6, 1.4, 5.0, 0.5, "Run started 2026-04-26 05:41 UTC",
              fill=C_ACCENT, size=12)
    _add_pill(s, 5.8, 1.4, 4.0, 0.5, "ETA ≈ 16:00 UTC (10 h wall)",
              fill=C_FOREST, size=12)
    _add_table(s, 0.6, 2.1, 8.5, 2.0, [
        ["r", "b_acc", "c_acc", "Δpp", "Verdict"],
        ["1", "54.0 %", "51.0 %", "−3.0", "reject (big overhaul)"],
        ["2", "64.0 %", "66.0 %", "+2.0", "ACCEPT"],
        ["3", "45.0 %", "40.0 %", "−5.0", "reject"],
    ])
    _add_bullets(s, 9.3, 2.1, 3.6, 4.5, [
        "Per-round wall: 45 / 50 / 73 min",
        "1 ACCEPT in 3 rounds so far",
        "Pace consistent with run #1",
        "7 rounds remaining",
        ("Slides will be updated once the run completes.", 0, C_MOSS),
    ], size=11)
    _add_textbox(s, 0.6, 6.5, 12.1, 0.6,
                 "Score on test acc + paired-accept rate (per pilot §8.4 option C). DO NOT re-introduce "
                 "best_val_acc > seed_batch_acc.",
                 size=12, italic=True, color=C_CHARCOAL)

    # --- 30. Section: novelty/weaknesses ---------------------------------
    make_section_slide(prs, "Novelty & weaknesses (honest)",
                       "What's actually new · what's still soft")

    # --- 31. Novelty ------------------------------------------------------
    s = make_content_slide(prs, "Novelty — the three real differentiators",
                           kicker="Honest self-assessment")
    _add_pill(s, 0.6, 1.4, 12.1, 0.5,
              "1.  Reflection-only multi-agent evolution (no search archive, no MCTS, no RL)",
              fill=C_FOREST, size=14)
    _add_bullets(s, 0.6, 2.0, 12.1, 1.4, [
        "ADAS / AFlow / GPTSwarm / MaAS / Puppeteer all rely on search or RL.",
        "Our controller is purely an in-context reflector over trajectory tapes on a frozen backbone.",
    ], size=13)
    _add_pill(s, 0.6, 3.5, 12.1, 0.5,
              "2.  Controller v2 — organization-designer framing + per-domain briefs + specialist personas",
              fill=C_FOREST, size=14)
    _add_bullets(s, 0.6, 4.1, 12.1, 1.4, [
        "Persona-authoring rules force domain vocabulary (gaap_analyst, differential_diagnostician, …).",
        "Briefs let the same backbone behave like a specialist on each domain without retraining.",
    ], size=13)
    _add_pill(s, 0.6, 5.6, 12.1, 0.5,
              "3.  Streaming paired-batch evolve mode — accept rate jumps from 1 / 9 to 4 / 10",
              fill=C_FOREST, size=14)
    _add_bullets(s, 0.6, 6.2, 12.1, 1.0, [
        "Eliminates cross-batch difficulty noise that was rejecting most architectural changes under legacy.",
    ], size=13)

    # --- 32. Weaknesses ---------------------------------------------------
    s = make_content_slide(prs, "Weaknesses — what we should NOT claim yet",
                           kicker="Honest self-assessment")
    _add_bullets(s, 0.6, 1.4, 12.1, 5.6, [
        ("v2 @ n=30 does NOT beat baselines on test in any of the three domains.", 0, C_ACCENT),
        ("±13 pp same-graph cross-run noise dominates the architectural effect at this n. The signal is below the noise floor.", 0, C_ACCENT),
        ("LLM-judge is same-family Qwen → self-bias risk on FinanceBench / AgentClinic.", 0, C_ACCENT),
        ("Single backbone (Qwen2.5-32B) and single judge.", 0, C_ACCENT),
        ("Streaming wall is 9-10 h × seed × domain → 3 × 3 sweep is ~90 h on a shared GPU.", 0, C_ACCENT),
        ("ADAS / MaAS / Puppeteer direct comparison NOT done yet (reviewer question #0).", 0, C_ACCENT),
        ("Concept-level anti-repeat is a soft constraint and may not fire — sanity_v3 round 2 was an exact repeat.", 0, C_ACCENT),
    ], size=14)

    # --- 33. What's next + timeline --------------------------------------
    s = make_content_slide(prs, "What's next + timeline",
                           kicker="Roadmap")
    _add_pill(s, 0.6, 1.5, 12.1, 0.5, "Critical path → EMNLP 2026 ARR (D-31, 2026-05-25)",
              fill=C_ACCENT, size=14)
    _add_bullets(s, 0.6, 2.1, 12.1, 4.9, [
        ("§5.2 multi-seed v2 streaming sweep (3 domains × 3 seeds, MEDIQ seed=1 running now). ~90 h.", 0, C_FOREST),
        "§5.3 Random-persona ablation (replace v2 personas with random text of same count).",
        "§5.4 Harness ablation (controller on/off, random topo, fixed topo).",
        "§5.5 Add a second backbone (Qwen3-72B + one of Claude / GPT API).",
        "§5.6 Replace LLM-judge with a different family (Claude Haiku 4.5 or GPT-4.1-mini).",
        ("§5.7 ADAS + (MaAS or Puppeteer or EvoMAC) direct baselines — non-negotiable for reviewer.", 0, C_ACCENT),
        "Failure-mode taxonomy (MAST-style), cost Pareto pipeline, judge-calibration experiment.",
    ], size=13)
    _add_textbox(s, 0.6, 7.0, 12.1, 0.4,
                 "Fallback venues: NeurIPS D&B (D-12, very aggressive) / EMNLP direct commitment (D-100) / NeurIPS workshops (summer).",
                 size=11, italic=True, color=C_CHARCOAL)

    # --- 34. Q&A ---------------------------------------------------------
    make_conclusion_slide(
        prs,
        "Thank you — Q&A",
        "Roadmap dashboard: references/roadmap.md · code: github.com/namam3gy (project repo)",
    )

    OUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT_PATH))
    print(f"wrote {OUT_PATH} ({len(prs.slides)} slides)")


if __name__ == "__main__":
    build()
